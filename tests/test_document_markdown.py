"""Step 2: document_markdown bridge (MarkItDown) — optional extra."""

from __future__ import annotations

import zipfile
from io import BytesIO
from pathlib import Path

import pytest

from adami_kernel.cortex import document_markdown as dm

pytest.importorskip("markitdown", reason="poetry install -E markitdown")


def _norm_escapes(markdown: str) -> str:
    """MarkItDown may escape underscores in plain runs."""
    return markdown.replace("\\_", "_")


def _write_minimal_docx(path: Path) -> None:
    document_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>MKStep2_DOCX_UNIQUE</w:t></w:r></w:p>
  </w:body>
</w:document>"""
    content_types = """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
    rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
    word_rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>"""
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("word/document.xml", document_xml)
        zf.writestr("word/_rels/document.xml.rels", word_rels)


def _write_sample_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    c.drawString(100, 750, "MKStep2_PDF_UNIQUE")
    c.showPage()
    c.save()


def _write_sample_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws["A1"] = "MKStep2_XLSX_UNIQUE"
    wb.save(str(path))


def _write_sample_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "MKStep2_PPTX_UNIQUE"
    prs.save(str(path))


@pytest.mark.asyncio
async def test_four_whitelisted_extensions_produce_markdown(tmp_path: Path) -> None:
    paths = {
        "pdf": tmp_path / "a.pdf",
        "docx": tmp_path / "b.docx",
        "xlsx": tmp_path / "c.xlsx",
        "pptx": tmp_path / "d.pptx",
    }
    _write_sample_pdf(paths["pdf"])
    _write_minimal_docx(paths["docx"])
    _write_sample_xlsx(paths["xlsx"])
    _write_sample_pptx(paths["pptx"])

    expected_substrings = {
        "pdf": "MKStep2_PDF_UNIQUE",
        "docx": "MKStep2_DOCX_UNIQUE",
        "xlsx": "MKStep2_XLSX_UNIQUE",
        "pptx": "MKStep2_PPTX_UNIQUE",
    }

    for key, p in paths.items():
        result = await dm.convert_document_path_to_markdown(p)
        assert isinstance(result, dm.DocumentMarkdownSuccess), (key, result)
        body = _norm_escapes(result.markdown)
        assert expected_substrings[key] in body, (key, result.markdown[:400])


@pytest.mark.asyncio
async def test_bad_pdf_bytes_returns_failure_not_uncaught(tmp_path: Path) -> None:
    p = tmp_path / "corrupt.pdf"
    # Valid PDF header but truncated body → PdfConverter raises (not plain-text fallback).
    p.write_bytes(b"%PDF-1.4\n1 0 obj<<>>\n")
    result = await dm.convert_document_path_to_markdown(p)
    assert isinstance(result, dm.DocumentMarkdownFailure)
    assert result.reason in (
        dm.DocumentMarkdownFailureReason.CONVERSION_FAILED,
        dm.DocumentMarkdownFailureReason.UNSUPPORTED_FORMAT,
        dm.DocumentMarkdownFailureReason.IO_ERROR,
    )


@pytest.mark.asyncio
async def test_disallowed_extension(tmp_path: Path) -> None:
    p = tmp_path / "nope.txt"
    p.write_text("hello", encoding="utf-8")
    result = await dm.convert_document_path_to_markdown(p)
    assert isinstance(result, dm.DocumentMarkdownFailure)
    assert result.reason == dm.DocumentMarkdownFailureReason.DISALLOWED_EXTENSION


@pytest.mark.asyncio
async def test_missing_path(tmp_path: Path) -> None:
    missing = tmp_path / "missing.docx"
    result = await dm.convert_document_path_to_markdown(missing)
    assert isinstance(result, dm.DocumentMarkdownFailure)
    assert result.reason == dm.DocumentMarkdownFailureReason.PATH_MISSING


@pytest.mark.asyncio
async def test_not_installed_branch(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4\n")
    monkeypatch.setattr(dm, "markitdown_import_spec", lambda: None)
    result = await dm.convert_document_path_to_markdown(p)
    assert isinstance(result, dm.DocumentMarkdownFailure)
    assert result.reason == dm.DocumentMarkdownFailureReason.NOT_INSTALLED


@pytest.mark.asyncio
async def test_truncation_sets_metadata(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    p = tmp_path / "dummy.pdf"
    p.write_bytes(b"%PDF-1.4\n")

    def _fake(_path_str: str) -> tuple[str, str | None]:
        return "Z" * 6000, None

    monkeypatch.setattr(dm, "markitdown_import_spec", lambda: object())
    monkeypatch.setattr(dm, "_run_markitdown_convert_path", _fake)
    result = await dm.convert_document_path_to_markdown(p, max_output_chars=100)
    assert isinstance(result, dm.DocumentMarkdownSuccess)
    assert result.meta.truncated is True
    assert result.meta.original_char_length == 6000
    assert len(result.markdown) == 100


@pytest.mark.asyncio
async def test_timeout_returns_timeout(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    def _slow(_path_str: str) -> tuple[str, str | None]:
        import time

        time.sleep(0.5)
        return "ok", None

    p = tmp_path / "wait.pdf"
    p.write_bytes(b"%PDF-1.4\n")
    monkeypatch.setattr(dm, "markitdown_import_spec", lambda: object())
    monkeypatch.setattr(dm, "_run_markitdown_convert_path", _slow)
    result = await dm.convert_document_path_to_markdown(p, timeout_s=0.05)
    assert isinstance(result, dm.DocumentMarkdownFailure)
    assert result.reason == dm.DocumentMarkdownFailureReason.TIMEOUT


@pytest.mark.asyncio
async def test_stream_docx_roundtrip(tmp_path: Path) -> None:
    docx = tmp_path / "streamed.docx"
    _write_minimal_docx(docx)
    raw = docx.read_bytes()
    from markitdown import StreamInfo

    bio = BytesIO(raw)
    si = StreamInfo(extension=".docx", filename="streamed.docx")
    result = await dm.convert_document_stream_to_markdown(bio, stream_info=si)
    assert isinstance(result, dm.DocumentMarkdownSuccess)
    assert "MKStep2_DOCX_UNIQUE" in _norm_escapes(result.markdown)


@pytest.mark.asyncio
async def test_stream_requires_whitelisted_extension() -> None:
    from markitdown import StreamInfo

    bio = BytesIO(b"hello")
    si = StreamInfo(extension=".txt", filename="x.txt")
    result = await dm.convert_document_stream_to_markdown(bio, stream_info=si)
    assert isinstance(result, dm.DocumentMarkdownFailure)
    assert result.reason == dm.DocumentMarkdownFailureReason.DISALLOWED_EXTENSION


@pytest.mark.parametrize(
    "name,expected_norm",
    [
        ("A.PDF", ".pdf"),
        ("x.DOCX", ".docx"),
        ("mix.PpTx", ".pptx"),
        ("sheet.XLSX", ".xlsx"),
    ],
)
def test_normalized_allowed_extension_case_insensitive(name: str, expected_norm: str) -> None:
    assert dm.normalized_allowed_extension(name) == expected_norm
